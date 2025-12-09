#!/usr/bin/env python3
"""
PPTX Merger Tool

Merges all PPTX files from a specified folder into a single presentation.
Usage: python tools/pptx_merger.py <folder_name>
"""

import argparse
import sys
from io import BytesIO
from pathlib import Path
from pptx import Presentation


def copy_slide(source_slide, dest_presentation):
    """
    Copy a slide from source presentation to destination presentation.
    
    Args:
        source_slide: Slide object from source presentation
        dest_presentation: Destination Presentation object
    """
    # Use blank layout for the new slide
    blank_slide_layout = dest_presentation.slide_layouts[6]
    dest_slide = dest_presentation.slides.add_slide(blank_slide_layout)
    
    # Copy all shapes from source to destination
    for shape in source_slide.shapes:
        try:
            # Handle pictures
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                # Convert bytes to BytesIO (file-like object) for add_picture()
                image_bytes = BytesIO(image.blob)
                dest_slide.shapes.add_picture(
                    image_bytes,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            # Handle text boxes and placeholders
            elif hasattr(shape, "text_frame") and shape.text_frame:
                # Create a text box with the same properties
                textbox = dest_slide.shapes.add_textbox(
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                text_frame = textbox.text_frame
                source_text_frame = shape.text_frame
                
                # Copy text content
                text_frame.text = source_text_frame.text
                
                # Copy paragraph formatting
                for i, source_paragraph in enumerate(source_text_frame.paragraphs):
                    if i < len(text_frame.paragraphs):
                        dest_paragraph = text_frame.paragraphs[i]
                        dest_paragraph.font.size = source_paragraph.font.size
                        if source_paragraph.font.bold is not None:
                            dest_paragraph.font.bold = source_paragraph.font.bold
                        if source_paragraph.font.italic is not None:
                            dest_paragraph.font.italic = source_paragraph.font.italic
            # Handle auto shapes
            elif hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
                new_shape = dest_slide.shapes.add_shape(
                    shape.auto_shape_type,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                if hasattr(shape, "text") and shape.text:
                    new_shape.text = shape.text
        except Exception as e:
            # If we can't copy a shape, log it but continue
            print(f"  Warning: Could not copy shape: {e}")
            continue


def merge_pptx_files(folder_name: str) -> None:
    """
    Merge all PPTX files from slideshows/{folder_name}/ into a single presentation.
    
    Args:
        folder_name: Name of the folder containing PPTX files to merge
    """
    # Get the script directory and project root
    script_dir = Path(__file__).parent
    project_root = script_dir.parent
    slideshows_dir = project_root / "slideshows"
    input_folder = slideshows_dir / folder_name
    output_file = slideshows_dir / f"{folder_name}.pptx"
    
    # Check if input folder exists
    if not input_folder.exists():
        print(f"Error: Folder '{input_folder}' does not exist.")
        sys.exit(1)
    
    if not input_folder.is_dir():
        print(f"Error: '{input_folder}' is not a directory.")
        sys.exit(1)
    
    # Find all PPTX files in the folder
    pptx_files = sorted(input_folder.glob("*.pptx"))
    
    if not pptx_files:
        print(f"Error: No PPTX files found in '{input_folder}'.")
        sys.exit(1)
    
    print(f"Found {len(pptx_files)} PPTX file(s) to merge:")
    for pptx_file in pptx_files:
        print(f"  - {pptx_file.name}")
    
    # Create a new presentation for the merged result
    merged_presentation = Presentation()
    
    # Remove the default slide if it exists
    if len(merged_presentation.slides) > 0:
        xml_slides = merged_presentation.slides._sldIdLst
        merged_presentation.part.drop_rel(xml_slides[0].rId)
        xml_slides.remove(xml_slides[0])
    
    # Process each PPTX file
    total_slides = 0
    for pptx_file in pptx_files:
        try:
            print(f"\nProcessing: {pptx_file.name}")
            source_presentation = Presentation(str(pptx_file))
            
            # Copy each slide from source to merged presentation
            for source_slide in source_presentation.slides:
                copy_slide(source_slide, merged_presentation)
                total_slides += 1
            
            print(f"  Added {len(source_presentation.slides)} slide(s)")
            
        except Exception as e:
            print(f"Error processing {pptx_file.name}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    # Save the merged presentation
    try:
        merged_presentation.save(str(output_file))
        print(f"\n{'='*50}")
        print(f"Success! Merged presentation saved to: {output_file}")
        print(f"Total slides: {total_slides}")
        print(f"{'='*50}")
    except Exception as e:
        print(f"Error saving merged presentation: {e}")
        sys.exit(1)


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description="Merge all PPTX files from a folder into a single presentation."
    )
    parser.add_argument(
        "folder_name",
        help="Name of the folder in slideshows/ containing PPTX files to merge"
    )
    
    args = parser.parse_args()
    merge_pptx_files(args.folder_name)


if __name__ == "__main__":
    main()
